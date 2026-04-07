"""Presentation action module for RabAI AutoClick.

Provides PowerPoint/presentation operations:
- PresentationCreateAction: Create new presentation
- PresentationAddSlideAction: Add a slide
- PresentationAddTextAction: Add text to slide
- PresentationAddImageAction: Add image to slide
- PresentationSaveAction: Save presentation
- PresentationReadAction: Read presentation content
- PresentationSlideCountAction: Get slide count
- PresentationDeleteSlideAction: Delete a slide
"""

import os
from typing import Any, Dict, List, Optional, Union

import sys
import os
_parent_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, _parent_dir)
from core.base_action import BaseAction, ActionResult

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class PresentationCreateAction(BaseAction):
    """Create new presentation."""
    action_type = "presentation_create"
    display_name = "创建演示文稿"
    description = "创建新的PowerPoint演示文稿"
    version = "1.0"

    def execute(
        self,
        context: Any,
        params: Dict[str, Any]
    ) -> ActionResult:
        """Execute create.

        Args:
            context: Execution context.
            params: Dict with file_path, width, height, layout.

        Returns:
            ActionResult indicating success.
        """
        file_path = params.get('file_path', '')
        width = params.get('width', 10)
        height = params.get('height', 7.5)
        layout = params.get('layout', 'title')

        valid, msg = self.validate_type(file_path, str, 'file_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        if not PPTX_AVAILABLE:
            return ActionResult(
                success=False,
                message="python-pptx未安装: pip install python-pptx"
            )

        try:
            resolved_path = context.resolve_value(file_path)
            resolved_width = context.resolve_value(width)
            resolved_height = context.resolve_value(height)
            resolved_layout = context.resolve_value(layout)

            prs = Presentation()
            prs.slide_width = Inches(resolved_width)
            prs.slide_height = Inches(resolved_height)

            # Add title slide
            layout_obj = prs.slide_layouts[6] if resolved_layout == 'blank' else prs.slide_layouts[0]
            prs.slides.add_slide(layout_obj)

            prs.save(resolved_path)

            return ActionResult(
                success=True,
                message=f"已创建演示文稿: {resolved_path}",
                data={'path': resolved_path}
            )
        except Exception as e:
            return ActionResult(
                success=False,
                message=f"创建演示文稿失败: {str(e)}"
            )

    def get_required_params(self) -> List[str]:
        return ['file_path']

    def get_optional_params(self) -> Dict[str, Any]:
        return {'width': 10, 'height': 7.5, 'layout': 'title'}


class PresentationAddSlideAction(BaseAction):
    """Add a slide."""
    action_type = "presentation_add_slide"
    display_name = "添加幻灯片"
    description = "向演示文稿添加新幻灯片"
    version = "1.0"

    def execute(
        self,
        context: Any,
        params: Dict[str, Any]
    ) -> ActionResult:
        """Execute add slide.

        Args:
            context: Execution context.
            params: Dict with file_path, layout, title.

        Returns:
            ActionResult indicating success.
        """
        file_path = params.get('file_path', '')
        layout = params.get('layout', 'title')
        title = params.get('title', '')

        valid, msg = self.validate_type(file_path, str, 'file_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        if not PPTX_AVAILABLE:
            return ActionResult(
                success=False,
                message="python-pptx未安装"
            )

        try:
            resolved_path = context.resolve_value(file_path)
            resolved_layout = context.resolve_value(layout)
            resolved_title = context.resolve_value(title) if title else ''

            if not os.path.exists(resolved_path):
                return ActionResult(
                    success=False,
                    message=f"文件不存在: {resolved_path}"
                )

            prs = Presentation(resolved_path)

            # Map layout name to index
            layout_map = {
                'title': 0, 'title_content': 1, 'section_header': 2,
                'two_content': 3, 'blank': 6, 'pic_tx': 7
            }
            layout_idx = layout_map.get(resolved_layout, 0)
            layout_obj = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]

            slide = prs.slides.add_slide(layout_obj)

            # Set title if provided
            if resolved_title:
                if slide.shapes.title:
                    slide.shapes.title.text = resolved_title

            prs.save(resolved_path)

            return ActionResult(
                success=True,
                message=f"已添加幻灯片: {resolved_path}",
                data={'slides': len(prs.slides), 'path': resolved_path}
            )
        except Exception as e:
            return ActionResult(
                success=False,
                message=f"添加幻灯片失败: {str(e)}"
            )

    def get_required_params(self) -> List[str]:
        return ['file_path']

    def get_optional_params(self) -> Dict[str, Any]:
        return {'layout': 'title', 'title': ''}


class PresentationAddTextAction(BaseAction):
    """Add text to slide."""
    action_type = "presentation_add_text"
    display_name = "添加文本"
    description = "向幻灯片添加文本框"
    version = "1.0"

    def execute(
        self,
        context: Any,
        params: Dict[str, Any]
    ) -> ActionResult:
        """Execute add text.

        Args:
            context: Execution context.
            params: Dict with file_path, slide_index, text, left, top, width, height.

        Returns:
            ActionResult indicating success.
        """
        file_path = params.get('file_path', '')
        slide_index = params.get('slide_index', 0)
        text = params.get('text', '')
        left = params.get('left', 1)
        top = params.get('top', 1)
        width = params.get('width', 5)
        height = params.get('height', 1)
        font_size = params.get('font_size', 18)
        bold = params.get('bold', False)
        color = params.get('color', '000000')

        valid, msg = self.validate_type(file_path, str, 'file_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        if not PPTX_AVAILABLE:
            return ActionResult(
                success=False,
                message="python-pptx未安装"
            )

        try:
            resolved_path = context.resolve_value(file_path)
            resolved_idx = context.resolve_value(slide_index)
            resolved_text = context.resolve_value(text)
            resolved_left = context.resolve_value(left)
            resolved_top = context.resolve_value(top)
            resolved_width = context.resolve_value(width)
            resolved_height = context.resolve_value(height)
            resolved_size = context.resolve_value(font_size)
            resolved_bold = context.resolve_value(bold)

            if not os.path.exists(resolved_path):
                return ActionResult(
                    success=False,
                    message=f"文件不存在: {resolved_path}"
                )

            prs = Presentation(resolved_path)

            if resolved_idx < 0 or resolved_idx >= len(prs.slides):
                return ActionResult(
                    success=False,
                    message=f"幻灯片索引无效: {resolved_idx}"
                )

            slide = prs.slides[resolved_idx]

            txbox = slide.shapes.add_textbox(
                Inches(resolved_left), Inches(resolved_top),
                Inches(resolved_width), Inches(resolved_height)
            )
            tf = txbox.text_frame
            tf.text = resolved_text

            # Format text
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(resolved_size)
                paragraph.font.bold = resolved_bold

                if color and color != '000000':
                    r = int(color[0:2], 16)
                    g = int(color[2:4], 16)
                    b = int(color[4:6], 16)
                    paragraph.font.color.rgb = RGBColor(r, g, b)

            prs.save(resolved_path)

            return ActionResult(
                success=True,
                message=f"已添加文本: {resolved_text[:50]}",
                data={'text': resolved_text[:100]}
            )
        except Exception as e:
            return ActionResult(
                success=False,
                message=f"添加文本失败: {str(e)}"
            )

    def get_required_params(self) -> List[str]:
        return ['file_path', 'slide_index', 'text']

    def get_optional_params(self) -> Dict[str, Any]:
        return {
            'left': 1, 'top': 1, 'width': 5, 'height': 1,
            'font_size': 18, 'bold': False, 'color': '000000'
        }


class PresentationAddImageAction(BaseAction):
    """Add image to slide."""
    action_type = "presentation_add_image"
    display_name = "添加图片"
    description = "向幻灯片添加图片"
    version = "1.0"

    def execute(
        self,
        context: Any,
        params: Dict[str, Any]
    ) -> ActionResult:
        """Execute add image.

        Args:
            context: Execution context.
            params: Dict with file_path, slide_index, image_path, left, top, width.

        Returns:
            ActionResult indicating success.
        """
        file_path = params.get('file_path', '')
        slide_index = params.get('slide_index', 0)
        image_path = params.get('image_path', '')
        left = params.get('left', 1)
        top = params.get('top', 1)
        width = params.get('width', 4)

        valid, msg = self.validate_type(file_path, str, 'file_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        valid, msg = self.validate_type(image_path, str, 'image_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        if not PPTX_AVAILABLE:
            return ActionResult(
                success=False,
                message="python-pptx未安装"
            )

        try:
            resolved_path = context.resolve_value(file_path)
            resolved_idx = context.resolve_value(slide_index)
            resolved_img = context.resolve_value(image_path)
            resolved_left = context.resolve_value(left)
            resolved_top = context.resolve_value(top)
            resolved_width = context.resolve_value(width)

            if not os.path.exists(resolved_path):
                return ActionResult(
                    success=False,
                    message=f"文件不存在: {resolved_path}"
                )

            if not os.path.exists(resolved_img):
                return ActionResult(
                    success=False,
                    message=f"图片不存在: {resolved_img}"
                )

            prs = Presentation(resolved_path)

            if resolved_idx < 0 or resolved_idx >= len(prs.slides):
                return ActionResult(
                    success=False,
                    message=f"幻灯片索引无效: {resolved_idx}"
                )

            slide = prs.slides[resolved_idx]
            slide.shapes.add_picture(
                resolved_img,
                Inches(resolved_left), Inches(resolved_top),
                width=Inches(resolved_width)
            )

            prs.save(resolved_path)

            return ActionResult(
                success=True,
                message=f"已添加图片: {resolved_img}",
                data={'image_path': resolved_img}
            )
        except Exception as e:
            return ActionResult(
                success=False,
                message=f"添加图片失败: {str(e)}"
            )

    def get_required_params(self) -> List[str]:
        return ['file_path', 'slide_index', 'image_path']

    def get_optional_params(self) -> Dict[str, Any]:
        return {'left': 1, 'top': 1, 'width': 4}


class PresentationSaveAction(BaseAction):
    """Save presentation."""
    action_type = "presentation_save"
    display_name = "保存演示文稿"
    description = "保存演示文稿"
    version = "1.0"

    def execute(
        self,
        context: Any,
        params: Dict[str, Any]
    ) -> ActionResult:
        """Execute save.

        Args:
            context: Execution context.
            params: Dict with file_path, output_path.

        Returns:
            ActionResult indicating success.
        """
        file_path = params.get('file_path', '')
        output_path = params.get('output_path', '')

        valid, msg = self.validate_type(file_path, str, 'file_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        if not PPTX_AVAILABLE:
            return ActionResult(
                success=False,
                message="python-pptx未安装"
            )

        try:
            resolved_path = context.resolve_value(file_path)
            resolved_output = context.resolve_value(output_path) if output_path else resolved_path

            if not os.path.exists(resolved_path):
                return ActionResult(
                    success=False,
                    message=f"文件不存在: {resolved_path}"
                )

            prs = Presentation(resolved_path)
            prs.save(resolved_output)

            return ActionResult(
                success=True,
                message=f"已保存: {resolved_output}",
                data={'path': resolved_output}
            )
        except Exception as e:
            return ActionResult(
                success=False,
                message=f"保存失败: {str(e)}"
            )

    def get_required_params(self) -> List[str]:
        return ['file_path']

    def get_optional_params(self) -> Dict[str, Any]:
        return {'output_path': ''}


class PresentationReadAction(BaseAction):
    """Read presentation content."""
    action_type = "presentation_read"
    display_name = "读取演示文稿"
    description = "读取演示文稿内容"
    version = "1.0"

    def execute(
        self,
        context: Any,
        params: Dict[str, Any]
    ) -> ActionResult:
        """Execute read.

        Args:
            context: Execution context.
            params: Dict with file_path, output_var.

        Returns:
            ActionResult with content.
        """
        file_path = params.get('file_path', '')
        output_var = params.get('output_var', 'presentation_content')

        valid, msg = self.validate_type(file_path, str, 'file_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        if not PPTX_AVAILABLE:
            return ActionResult(
                success=False,
                message="python-pptx未安装"
            )

        try:
            resolved_path = context.resolve_value(file_path)

            if not os.path.exists(resolved_path):
                return ActionResult(
                    success=False,
                    message=f"文件不存在: {resolved_path}"
                )

            prs = Presentation(resolved_path)
            slides_data = []

            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                slides_data.append({
                    'index': i,
                    'texts': slide_texts
                })

            context.set(output_var, slides_data)

            return ActionResult(
                success=True,
                message=f"已读取 {len(slides_data)} 张幻灯片",
                data={'slides': len(slides_data), 'output_var': output_var}
            )
        except Exception as e:
            return ActionResult(
                success=False,
                message=f"读取演示文稿失败: {str(e)}"
            )

    def get_required_params(self) -> List[str]:
        return ['file_path']

    def get_optional_params(self) -> Dict[str, Any]:
        return {'output_var': 'presentation_content'}


class PresentationSlideCountAction(BaseAction):
    """Get slide count."""
    action_type = "presentation_slide_count"
    display_name = "获取幻灯片数量"
    description = "获取演示文稿的幻灯片数量"
    version = "1.0"

    def execute(
        self,
        context: Any,
        params: Dict[str, Any]
    ) -> ActionResult:
        """Execute count.

        Args:
            context: Execution context.
            params: Dict with file_path, output_var.

        Returns:
            ActionResult with slide count.
        """
        file_path = params.get('file_path', '')
        output_var = params.get('output_var', 'slide_count')

        valid, msg = self.validate_type(file_path, str, 'file_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        if not PPTX_AVAILABLE:
            return ActionResult(
                success=False,
                message="python-pptx未安装"
            )

        try:
            resolved_path = context.resolve_value(file_path)

            if not os.path.exists(resolved_path):
                return ActionResult(
                    success=False,
                    message=f"文件不存在: {resolved_path}"
                )

            prs = Presentation(resolved_path)
            count = len(prs.slides)

            context.set(output_var, count)

            return ActionResult(
                success=True,
                message=f"幻灯片数量: {count}",
                data={'count': count, 'output_var': output_var}
            )
        except Exception as e:
            return ActionResult(
                success=False,
                message=f"获取幻灯片数量失败: {str(e)}"
            )

    def get_required_params(self) -> List[str]:
        return ['file_path']

    def get_optional_params(self) -> Dict[str, Any]:
        return {'output_var': 'slide_count'}


class PresentationDeleteSlideAction(BaseAction):
    """Delete a slide."""
    action_type = "presentation_delete_slide"
    display_name = "删除幻灯片"
    description = "删除指定幻灯片"
    version = "1.0"

    def execute(
        self,
        context: Any,
        params: Dict[str, Any]
    ) -> ActionResult:
        """Execute delete.

        Args:
            context: Execution context.
            params: Dict with file_path, slide_index.

        Returns:
            ActionResult indicating success.
        """
        file_path = params.get('file_path', '')
        slide_index = params.get('slide_index', 0)

        valid, msg = self.validate_type(file_path, str, 'file_path')
        if not valid:
            return ActionResult(success=False, message=msg)

        if not PPTX_AVAILABLE:
            return ActionResult(
                success=False,
                message="python-pptx未安装"
            )

        try:
            resolved_path = context.resolve_value(file_path)
            resolved_idx = context.resolve_value(slide_index)

            if not os.path.exists(resolved_path):
                return ActionResult(
                    success=False,
                    message=f"文件不存在: {resolved_path}"
                )

            prs = Presentation(resolved_path)

            if resolved_idx < 0 or resolved_idx >= len(prs.slides):
                return ActionResult(
                    success=False,
                    message=f"幻灯片索引无效: {resolved_idx}"
                )

            # Get the slide's rId
            slide_id = prs.slides._sldIdLst[resolved_idx].rId
            prs.part.drop_rel(slide_id)
            del prs.slides._sldIdLst[resolved_idx]

            prs.save(resolved_path)

            return ActionResult(
                success=True,
                message=f"已删除第 {resolved_idx + 1} 张幻灯片",
                data={'deleted_index': resolved_idx}
            )
        except Exception as e:
            return ActionResult(
                success=False,
                message=f"删除幻灯片失败: {str(e)}"
            )

    def get_required_params(self) -> List[str]:
        return ['file_path', 'slide_index']

    def get_optional_params(self) -> Dict[str, Any]:
        return {}
